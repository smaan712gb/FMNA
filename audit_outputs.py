"""
Comprehensive Output Quality Audit Script
Analyzes all output files in fe_results for placeholders, blank data, and quality issues
"""

import os
import sys
from pathlib import Path
import json
from datetime import datetime
import re

# Handle different file types
try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("Warning: openpyxl not installed. Excel files will be skipped.")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not installed. DOCX files will be skipped.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. PPTX files will be skipped.")

class OutputAuditor:
    """Audits output files for quality issues"""
    
    # Patterns that indicate placeholder or incomplete content
    PLACEHOLDER_PATTERNS = [
        r'\[.*?\]',  # [placeholder]
        r'<.*?>',    # <placeholder>
        r'{.*?}',    # {placeholder}
        r'TODO',
        r'FIXME',
        r'XXX',
        r'TBD',
        r'TBA',
        r'N/A',
        r'placeholder',
        r'dummy',
        r'test data',
        r'sample',
        r'example\.com',
        r'lorem ipsum',
        r'your .* here',
        r'insert .* here',
        r'to be determined',
        r'pending',
        r'undefined',
        r'null',
        r'None',
        r'#REF!',
        r'#VALUE!',
        r'#DIV/0!',
        r'#N/A',
        r'0\.00%',  # Potentially placeholder percentages
        r'\$0\.00',  # Potentially placeholder currency
        r'^\s*$',    # Empty cells/content
    ]
    
    def __init__(self, results_dir='fe_results'):
        self.results_dir = Path(results_dir)
        self.issues = []
        self.file_stats = {}
        
    def audit_all_files(self):
        """Main audit function"""
        print(f"\n{'='*80}")
        print(f"OUTPUT QUALITY AUDIT - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        print(f"{'='*80}\n")
        
        if not self.results_dir.exists():
            print(f"❌ Error: Directory {self.results_dir} does not exist")
            return
            
        files = list(self.results_dir.glob('*'))
        if not files:
            print(f"❌ Error: No files found in {self.results_dir}")
            return
            
        print(f"📁 Auditing {len(files)} files in {self.results_dir}\n")
        
        for filepath in sorted(files):
            if filepath.is_file():
                self.audit_file(filepath)
        
        self.generate_report()
        
    def audit_file(self, filepath: Path):
        """Audit a single file based on its type"""
        print(f"🔍 Auditing: {filepath.name}")
        
        file_issues = []
        file_stats = {
            'name': filepath.name,
            'size': filepath.stat().st_size,
            'type': filepath.suffix,
            'issues': file_issues
        }
        
        try:
            if filepath.suffix == '.xlsx' and EXCEL_AVAILABLE:
                file_issues.extend(self.audit_excel(filepath))
            elif filepath.suffix == '.docx' and DOCX_AVAILABLE:
                file_issues.extend(self.audit_docx(filepath))
            elif filepath.suffix == '.pptx' and PPTX_AVAILABLE:
                file_issues.extend(self.audit_pptx(filepath))
            elif filepath.suffix == '.html':
                file_issues.extend(self.audit_html(filepath))
            else:
                file_issues.append(f"Unsupported file type or missing library: {filepath.suffix}")
                
        except Exception as e:
            file_issues.append(f"Error auditing file: {str(e)}")
            
        self.file_stats[filepath.name] = file_stats
        
        if file_issues:
            print(f"  ⚠️  Found {len(file_issues)} issues")
            for issue in file_issues[:3]:  # Show first 3 issues
                print(f"      - {issue}")
            if len(file_issues) > 3:
                print(f"      ... and {len(file_issues) - 3} more")
        else:
            print(f"  ✅ No issues found")
        print()
        
    def check_for_placeholders(self, text: str, location: str) -> list:
        """Check text for placeholder patterns"""
        issues = []
        if not text or not isinstance(text, str):
            return issues
            
        for pattern in self.PLACEHOLDER_PATTERNS:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                # Get context around the match
                start = max(0, match.start() - 30)
                end = min(len(text), match.end() + 30)
                context = text[start:end].replace('\n', ' ').strip()
                
                issues.append(f"{location}: '{match.group()}' in context: ...{context}...")
                
        return issues
        
    def audit_excel(self, filepath: Path) -> list:
        """Audit Excel file"""
        issues = []
        try:
            wb = load_workbook(filepath, data_only=True)
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # Check for empty sheets
                if sheet.max_row <= 1:
                    issues.append(f"Sheet '{sheet_name}': Empty or only headers")
                    continue
                
                # Sample cells from the sheet
                for row_idx, row in enumerate(sheet.iter_rows(min_row=2, max_row=min(100, sheet.max_row), values_only=True), start=2):
                    for col_idx, cell_value in enumerate(row, start=1):
                        if cell_value is not None:
                            cell_text = str(cell_value)
                            cell_ref = f"Sheet '{sheet_name}', Cell {chr(64 + col_idx)}{row_idx}"
                            issues.extend(self.check_for_placeholders(cell_text, cell_ref))
                            
            wb.close()
            
        except Exception as e:
            issues.append(f"Error reading Excel: {str(e)}")
            
        return issues
        
    def audit_docx(self, filepath: Path) -> list:
        """Audit Word document"""
        issues = []
        try:
            doc = Document(filepath)
            
            # Check if document is empty
            if len(doc.paragraphs) == 0:
                issues.append("Document has no paragraphs")
                return issues
                
            # Check paragraphs
            for para_idx, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    issues.extend(self.check_for_placeholders(
                        para.text, 
                        f"Paragraph {para_idx + 1}"
                    ))
                    
            # Check tables
            for table_idx, table in enumerate(doc.tables):
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            issues.extend(self.check_for_placeholders(
                                cell.text,
                                f"Table {table_idx + 1}, Row {row_idx + 1}, Cell {cell_idx + 1}"
                            ))
                            
        except Exception as e:
            issues.append(f"Error reading DOCX: {str(e)}")
            
        return issues
        
    def audit_pptx(self, filepath: Path) -> list:
        """Audit PowerPoint presentation"""
        issues = []
        try:
            prs = Presentation(filepath)
            
            if len(prs.slides) == 0:
                issues.append("Presentation has no slides")
                return issues
                
            for slide_idx, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        issues.extend(self.check_for_placeholders(
                            shape.text,
                            f"Slide {slide_idx}"
                        ))
                        
        except Exception as e:
            issues.append(f"Error reading PPTX: {str(e)}")
            
        return issues
        
    def audit_html(self, filepath: Path) -> list:
        """Audit HTML file"""
        issues = []
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read()
                
            # Remove HTML tags for text analysis
            text_only = re.sub(r'<[^>]+>', ' ', content)
            issues.extend(self.check_for_placeholders(text_only, filepath.name))
            
        except Exception as e:
            issues.append(f"Error reading HTML: {str(e)}")
            
        return issues
        
    def generate_report(self):
        """Generate comprehensive audit report"""
        print(f"\n{'='*80}")
        print("AUDIT REPORT SUMMARY")
        print(f"{'='*80}\n")
        
        total_files = len(self.file_stats)
        files_with_issues = sum(1 for stats in self.file_stats.values() if stats['issues'])
        total_issues = sum(len(stats['issues']) for stats in self.file_stats.values())
        
        print(f"📊 Overall Statistics:")
        print(f"   Total files audited: {total_files}")
        print(f"   Files with issues: {files_with_issues}")
        print(f"   Total issues found: {total_issues}")
        print()
        
        if files_with_issues > 0:
            print(f"⚠️  FILES REQUIRING ATTENTION:\n")
            for filename, stats in sorted(self.file_stats.items()):
                if stats['issues']:
                    print(f"  📄 {filename}")
                    print(f"     Issues: {len(stats['issues'])}")
                    print(f"     Size: {stats['size']:,} bytes")
                    for issue in stats['issues'][:5]:  # Show first 5
                        print(f"       • {issue}")
                    if len(stats['issues']) > 5:
                        print(f"       ... and {len(stats['issues']) - 5} more issues")
                    print()
        else:
            print("✅ All files passed quality audit!\n")
            
        # Save detailed report
        report_path = self.results_dir / f'audit_report_{datetime.now().strftime("%Y%m%d_%H%M")}.json'
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(self.file_stats, f, indent=2, default=str)
        print(f"💾 Detailed report saved to: {report_path}")
        
        return total_issues == 0

def main():
    """Main execution"""
    auditor = OutputAuditor('fe_results')
    auditor.audit_all_files()
    
if __name__ == '__main__':
    main()
