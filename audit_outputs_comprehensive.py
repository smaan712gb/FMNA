"""
Comprehensive Output Quality Audit Script - Professional M&A Grade
Detects: Empty sheets, placeholders, generic data, missing real data
"""

import os
import sys
from pathlib import Path
import json
from datetime import datetime
import re
from collections import defaultdict

# Handle different file types
try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("⚠️  WARNING: openpyxl not installed. Excel files will be skipped.")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("⚠️  WARNING: python-docx not installed. DOCX files will be skipped.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  WARNING: python-pptx not installed. PPTX files will be skipped.")

class ComprehensiveOutputAuditor:
    """Professional M&A grade output auditor"""
    
    # Critical placeholder patterns
    CRITICAL_PLACEHOLDERS = [
        r'\[.*?\]',
        r'TODO',
        r'FIXME',
        r'XXX',
        r'TBD',
        r'TBA',
        r'placeholder',
        r'dummy',
        r'test data',
        r'sample',
        r'example\.com',
        r'lorem ipsum',
        r'your .* here',
        r'insert .* here',
        r'to be determined',
        r'pending',
        r'undefined',
        r'#REF!',
        r'#VALUE!',
        r'#DIV/0!',
        r'#N/A',
    ]
    
    # Generic/template indicators
    GENERIC_INDICATORS = [
        r'Company Name',
        r'Target Company',
        r'Acquirer',
        r'Sample Corp',
        r'Example Inc',
        r'ABC Company',
        r'XYZ Corp',
        r'Client Name',
    ]
    
    def __init__(self, results_dir='fe_results', company_ticker='NVDA'):
        self.results_dir = Path(results_dir)
        self.company_ticker = company_ticker.upper()
        self.issues = defaultdict(list)
        self.file_stats = {}
        self.critical_issues = []
        
    def audit_all_files(self):
        """Main audit function"""
        print(f"\n{'='*100}")
        print(f"COMPREHENSIVE M&A OUTPUT AUDIT - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        print(f"Company: {self.company_ticker}")
        print(f"{'='*100}\n")
        
        if not self.results_dir.exists():
            print(f"❌ CRITICAL: Directory {self.results_dir} does not exist")
            return False
            
        files = list(self.results_dir.glob('*'))
        if not files:
            print(f"❌ CRITICAL: No files found in {self.results_dir}")
            return False
            
        print(f"📁 Auditing {len(files)} files in {self.results_dir}\n")
        
        for filepath in sorted(files):
            if filepath.is_file() and not filepath.name.startswith('audit_report'):
                self.audit_file(filepath)
        
        return self.generate_report()
        
    def audit_file(self, filepath: Path):
        """Audit a single file based on its type"""
        print(f"🔍 Auditing: {filepath.name}")
        
        file_issues = []
        
        try:
            if filepath.suffix == '.xlsx' and EXCEL_AVAILABLE:
                file_issues = self.audit_excel_comprehensive(filepath)
            elif filepath.suffix == '.docx' and DOCX_AVAILABLE:
                file_issues = self.audit_docx_comprehensive(filepath)
            elif filepath.suffix == '.pptx' and PPTX_AVAILABLE:
                file_issues = self.audit_pptx_comprehensive(filepath)
            elif filepath.suffix == '.html':
                file_issues = self.audit_html_comprehensive(filepath)
            else:
                file_issues.append({
                    'severity': 'ERROR',
                    'message': f"Unsupported file type or missing library: {filepath.suffix}"
                })
                
        except Exception as e:
            file_issues.append({
                'severity': 'CRITICAL',
                'message': f"Error auditing file: {str(e)}"
            })
            
        self.issues[filepath.name] = file_issues
        
        critical_count = sum(1 for i in file_issues if i.get('severity') in ['CRITICAL', 'ERROR'])
        warning_count = sum(1 for i in file_issues if i.get('severity') == 'WARNING')
        
        if critical_count > 0:
            print(f"  ❌ CRITICAL: {critical_count} critical issues found")
            self.critical_issues.extend(file_issues)
        elif warning_count > 0:
            print(f"  ⚠️  {warning_count} warnings found")
        else:
            print(f"  ✅ Passed quality checks")
        print()
        
    def audit_excel_comprehensive(self, filepath: Path) -> list:
        """Comprehensive Excel audit"""
        issues = []
        
        try:
            wb = load_workbook(filepath, data_only=True)
            
            # Check each sheet
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_issues = self.audit_excel_sheet(sheet, sheet_name)
                issues.extend(sheet_issues)
                    
            wb.close()
            
        except Exception as e:
            issues.append({
                'severity': 'CRITICAL',
                'location': 'Excel File',
                'message': f"Cannot read Excel file: {str(e)}"
            })
            
        return issues
        
    def audit_excel_sheet(self, sheet, sheet_name: str) -> list:
        """Audit individual Excel sheet"""
        issues = []
        
        # Check if sheet is completely empty
        if sheet.max_row == 0 or sheet.max_column == 0:
            issues.append({
                'severity': 'CRITICAL',
                'location': f"Sheet: {sheet_name}",
                'message': 'Sheet is completely empty - no data at all'
            })
            return issues
            
        # Check if sheet only has headers (1 row)
        if sheet.max_row == 1:
            issues.append({
                'severity': 'CRITICAL',
                'location': f"Sheet: {sheet_name}",
                'message': 'Sheet only contains headers - no actual data'
            })
            return issues
            
        # Count non-empty cells
        non_empty_count = 0
        total_cells_checked = 0
        zero_count = 0
        null_like_count = 0
        
        # Check meaningful data (skip first row as headers)
        for row_idx, row in enumerate(sheet.iter_rows(min_row=2, max_row=min(1000, sheet.max_row)), start=2):
            for col_idx, cell in enumerate(row, start=1):
                total_cells_checked += 1
                
                if cell.value is not None:
                    non_empty_count += 1
                    cell_str = str(cell.value).strip()
                    
                    # Check for zeros (might indicate missing data)
                    if cell_str in ['0', '0.0', '0.00']:
                        zero_count += 1
                    
                    # Check for null-like values
                    if cell_str.lower() in ['none', 'null', 'n/a', 'na', '#n/a', '-']:
                        null_like_count += 1
                        
                    # Check for placeholders in this cell
                    for pattern in self.CRITICAL_PLACEHOLDERS:
                        if re.search(pattern, cell_str, re.IGNORECASE):
                            cell_ref = f"{chr(64 + col_idx) if col_idx <= 26 else 'Col' + str(col_idx)}{row_idx}"
                            issues.append({
                                'severity': 'CRITICAL',
                                'location': f"Sheet: {sheet_name}, Cell: {cell_ref}",
                                'message': f"Placeholder found: '{cell_str[:50]}'"
                            })
                    
                    # Check for generic company names
                    for pattern in self.GENERIC_INDICATORS:
                        if re.search(pattern, cell_str, re.IGNORECASE):
                            cell_ref = f"{chr(64 + col_idx) if col_idx <= 26 else 'Col' + str(col_idx)}{row_idx}"
                            issues.append({
                                'severity': 'ERROR',
                                'location': f"Sheet: {sheet_name}, Cell: {cell_ref}",
                                'message': f"Generic placeholder company name: '{cell_str[:50]}'"
                            })
        
        # Analyze data density
        if total_cells_checked > 0:
            empty_percentage = ((total_cells_checked - non_empty_count) / total_cells_checked) * 100
            zero_percentage = (zero_count / non_empty_count * 100) if non_empty_count > 0 else 0
            
            # Flag sheets with mostly empty cells
            if empty_percentage > 80:
                issues.append({
                    'severity': 'CRITICAL',
                    'location': f"Sheet: {sheet_name}",
                    'message': f'Sheet is {empty_percentage:.1f}% empty - likely placeholder/template'
                })
            elif empty_percentage > 50:
                issues.append({
                    'severity': 'WARNING',
                    'location': f"Sheet: {sheet_name}",
                    'message': f'Sheet is {empty_percentage:.1f}% empty - check data completeness'
                })
                
            # Flag sheets with too many zeros
            if zero_percentage > 50:
                issues.append({
                    'severity': 'ERROR',
                    'location': f"Sheet: {sheet_name}",
                    'message': f'{zero_percentage:.1f}% of cells are zeros - likely missing real data'
                })
                
            # Flag sheets with many null-like values
            if null_like_count > non_empty_count * 0.3:
                issues.append({
                    'severity': 'ERROR',
                    'location': f"Sheet: {sheet_name}",
                    'message': f'{null_like_count} cells contain null-like values (N/A, None, etc.)'
                })
        
        # Check for company-specific data
        if not self.check_company_specific_data_excel(sheet, sheet_name):
            issues.append({
                'severity': 'WARNING',
                'location': f"Sheet: {sheet_name}",
                'message': f'No {self.company_ticker}-specific data found - may be generic template'
            })
            
        return issues
        
    def check_company_specific_data_excel(self, sheet, sheet_name: str) -> bool:
        """Check if Excel sheet contains company-specific data"""
        company_mentions = 0
        
        for row in sheet.iter_rows(max_row=100):
            for cell in row:
                if cell.value:
                    cell_str = str(cell.value).upper()
                    if self.company_ticker in cell_str:
                        company_mentions += 1
                    # Also check for common company names
                    if self.company_ticker == 'NVDA' and 'NVIDIA' in cell_str:
                        company_mentions += 1
                        
        return company_mentions > 0
        
    def audit_docx_comprehensive(self, filepath: Path) -> list:
        """Comprehensive Word document audit"""
        issues = []
        
        try:
            doc = Document(filepath)
            
            # Check if document is empty
            if len(doc.paragraphs) == 0 and len(doc.tables) == 0:
                issues.append({
                    'severity': 'CRITICAL',
                    'location': 'Document',
                    'message': 'Document is completely empty'
                })
                return issues
            
            # Count meaningful content
            text_content = []
            
            # Check paragraphs
            for para_idx, para in enumerate(doc.paragraphs):
                text = para.text.strip()
                if text:
                    text_content.append(text)
                    
                    # Check for placeholders
                    for pattern in self.CRITICAL_PLACEHOLDERS:
                        if re.search(pattern, text, re.IGNORECASE):
                            issues.append({
                                'severity': 'CRITICAL',
                                'location': f"Paragraph {para_idx + 1}",
                                'message': f"Placeholder: '{text[:100]}'"
                            })
                    
                    # Check for generic company names
                    for pattern in self.GENERIC_INDICATORS:
                        if re.search(pattern, text, re.IGNORECASE):
                            issues.append({
                                'severity': 'ERROR',
                                'location': f"Paragraph {para_idx + 1}",
                                'message': f"Generic company reference: '{text[:100]}'"
                            })
            
            # Check tables
            for table_idx, table in enumerate(doc.tables):
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        if text:
                            text_content.append(text)
                            
                            # Look for "None required" or similar
                            if re.search(r'\bnone\s+required\b', text, re.IGNORECASE):
                                issues.append({
                                    'severity': 'WARNING',
                                    'location': f"Table {table_idx + 1}, Row {row_idx + 1}, Cell {cell_idx + 1}",
                                    'message': f"Generic text 'None required': '{text[:100]}'"
                                })
                            
                            # Check for other placeholders
                            for pattern in self.CRITICAL_PLACEHOLDERS:
                                if re.search(pattern, text, re.IGNORECASE):
                                    issues.append({
                                        'severity': 'CRITICAL',
                                        'location': f"Table {table_idx + 1}, Row {row_idx + 1}, Cell {cell_idx + 1}",
                                        'message': f"Placeholder: '{text[:100]}'"
                                    })
            
            # Check overall content quality
            all_text = ' '.join(text_content)
            
            if len(all_text) < 500:
                issues.append({
                    'severity': 'ERROR',
                    'location': 'Document',
                    'message': f'Document has very little content ({len(all_text)} characters) - likely incomplete'
                })
            
            # Check for company-specific content
            company_mentions = all_text.upper().count(self.company_ticker)
            if self.company_ticker == 'NVDA':
                company_mentions += all_text.upper().count('NVIDIA')
                
            if company_mentions == 0:
                issues.append({
                    'severity': 'ERROR',
                    'location': 'Document',
                    'message': f'No mention of {self.company_ticker} found - may be generic template'
                })
            elif company_mentions < 3:
                issues.append({
                    'severity': 'WARNING',
                    'location': 'Document',
                    'message': f'Only {company_mentions} mentions of {self.company_ticker} - check specificity'
                })
                
        except Exception as e:
            issues.append({
                'severity': 'CRITICAL',
                'location': 'Document',
                'message': f"Cannot read DOCX: {str(e)}"
            })
            
        return issues
        
    def audit_pptx_comprehensive(self, filepath: Path) -> list:
        """Comprehensive PowerPoint audit"""
        issues = []
        
        try:
            prs = Presentation(filepath)
            
            if len(prs.slides) == 0:
                issues.append({
                    'severity': 'CRITICAL',
                    'location': 'Presentation',
                    'message': 'Presentation has no slides'
                })
                return issues
            
            # Check each slide
            all_text = []
            
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_text.append(text)
                        all_text.append(text)
                        
                        # Check for placeholders
                        for pattern in self.CRITICAL_PLACEHOLDERS:
                            if re.search(pattern, text, re.IGNORECASE):
                                issues.append({
                                    'severity': 'CRITICAL',
                                    'location': f"Slide {slide_idx}",
                                    'message': f"Placeholder: '{text[:100]}'"
                                })
                        
                        # Check for generic company names
                        for pattern in self.GENERIC_INDICATORS:
                            if re.search(pattern, text, re.IGNORECASE):
                                issues.append({
                                    'severity': 'ERROR',
                                    'location': f"Slide {slide_idx}",
                                    'message': f"Generic company: '{text[:100]}'"
                                })
                
                # Check if slide is mostly empty
                if len(' '.join(slide_text)) < 50:
                    issues.append({
                        'severity': 'WARNING',
                        'location': f"Slide {slide_idx}",
                        'message': 'Slide has minimal content'
                    })
            
            # Check company-specific content
            combined_text = ' '.join(all_text).upper()
            company_mentions = combined_text.count(self.company_ticker)
            
            if company_mentions == 0:
                issues.append({
                    'severity': 'ERROR',
                    'location': 'Presentation',
                    'message': f'No mention of {self.company_ticker}'
                })
                
        except Exception as e:
            issues.append({
                'severity': 'CRITICAL',
                'location': 'Presentation',
                'message': f"Cannot read PPTX: {str(e)}"
            })
            
        return issues
        
    def audit_html_comprehensive(self, filepath: Path) -> list:
        """Comprehensive HTML audit - ignore JavaScript"""
        issues = []
        
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Extract text content (visible text, not JavaScript)
            # Remove script tags entirely
            content_no_script = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL)
            
            # Remove other tags
            text_only = re.sub(r'<[^>]+>', ' ', content_no_script)
            
            # Check for meaningful content
            if len(text_only.strip()) < 100:
                issues.append({
                    'severity': 'ERROR',
                    'location': 'HTML',
                    'message': 'Very little visible text content'
                })
            
            # Check for placeholders in visible text only
            for pattern in self.CRITICAL_PLACEHOLDERS:
                matches = list(re.finditer(pattern, text_only, re.IGNORECASE))
                # Only report if more than 5 matches (to avoid JavaScript false positives)
                if len(matches) > 5:
                    issues.append({
                        'severity': 'WARNING',
                        'location': 'HTML Content',
                        'message': f"Multiple '{matches[0].group()}' patterns found"
                    })
                    break
                    
        except Exception as e:
            issues.append({
                'severity': 'ERROR',
                'location': 'HTML',
                'message': f"Cannot read HTML: {str(e)}"
            })
            
        return issues
        
    def generate_report(self):
        """Generate comprehensive audit report"""
        print(f"\n{'='*100}")
        print("COMPREHENSIVE AUDIT REPORT")
        print(f"{'='*100}\n")
        
        total_files = len(self.issues)
        files_with_critical = sum(1 for issues in self.issues.values() 
                                  if any(i.get('severity') == 'CRITICAL' for i in issues))
        files_with_errors = sum(1 for issues in self.issues.values() 
                               if any(i.get('severity') == 'ERROR' for i in issues))
        files_with_warnings = sum(1 for issues in self.issues.values() 
                                 if any(i.get('severity') == 'WARNING' for i in issues))
        
        total_critical = sum(sum(1 for i in issues if i.get('severity') == 'CRITICAL') 
                           for issues in self.issues.values())
        total_errors = sum(sum(1 for i in issues if i.get('severity') == 'ERROR') 
                          for issues in self.issues.values())
        total_warnings = sum(sum(1 for i in issues if i.get('severity') == 'WARNING') 
                            for issues in self.issues.values())
        
        print(f"📊 OVERALL STATISTICS:")
        print(f"   Total files audited: {total_files}")
        print(f"   Files with CRITICAL issues: {files_with_critical}")
        print(f"   Files with ERROR issues: {files_with_errors}")
        print(f"   Files with WARNINGS: {files_with_warnings}")
        print()
        print(f"   Total CRITICAL issues: {total_critical}")
        print(f"   Total ERROR issues: {total_errors}")
        print(f"   Total WARNINGS: {total_warnings}")
        print()
        
        # Determine pass/fail
        passed = total_critical == 0 and total_errors == 0
        
        if passed:
            print("✅ ALL FILES PASSED M&A QUALITY STANDARDS")
            print("   Minor warnings may exist but no critical issues found.\n")
        else:
            print("❌ OUTPUTS FAILED M&A QUALITY STANDARDS")
            print("   Critical issues and errors must be fixed before client delivery.\n")
        
        # Show detailed issues
        if files_with_critical > 0 or files_with_errors > 0:
            print(f"{'='*100}")
            print("DETAILED ISSUES REQUIRING IMMEDIATE ATTENTION:")
            print(f"{'='*100}\n")
            
            for filename, issues in sorted(self.issues.items()):
                critical_issues = [i for i in issues if i.get('severity') == 'CRITICAL']
                error_issues = [i for i in issues if i.get('severity') == 'ERROR']
                
                if critical_issues or error_issues:
                    print(f"📄 {filename}")
                    print(f"   CRITICAL: {len(critical_issues)}, ERRORS: {len(error_issues)}")
                    
                    # Show critical issues
                    for issue in critical_issues[:10]:
                        print(f"   ❌ CRITICAL [{issue.get('location', 'N/A')}]: {issue.get('message', 'N/A')}")
                    
                    # Show error issues
                    for issue in error_issues[:10]:
                        print(f"   ⚠️  ERROR [{issue.get('location', 'N/A')}]: {issue.get('message', 'N/A')}")
                    
                    remaining = len(critical_issues) + len(error_issues) - 20
                    if remaining > 0:
                        print(f"   ... and {remaining} more issues")
                    print()
        
        # Save detailed report
        report_path = self.results_dir / f'comprehensive_audit_report_{datetime.now().strftime("%Y%m%d_%H%M")}.json'
        report_data = {
            'audit_time': datetime.now().isoformat(),
            'company_ticker': self.company_ticker,
            'statistics': {
                'total_files': total_files,
                'files_with_critical': files_with_critical,
                'files_with_errors': files_with_errors,
                'files_with_warnings': files_with_warnings,
                'total_critical': total_critical,
                'total_errors': total_errors,
                'total_warnings': total_warnings,
                'passed': passed
            },
            'issues': {filename: issues for filename, issues in self.issues.items()}
        }
        
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report_data, f, indent=2)
        
        print(f"💾 Detailed report saved to: {report_path}")
        print(f"\n{'='*100}\n")
        
        return passed

def main():
    """Main execution"""
    import argparse
    parser = argparse.ArgumentParser(description='Comprehensive M&A Output Quality Audit')
    parser.add_argument('--dir', default='fe_results', help='Results directory')
    parser.add_argument('--ticker', default='NVDA', help='Company ticker symbol')
    
    args = parser.parse_args()
    
    auditor = ComprehensiveOutputAuditor(args.dir, args.ticker)
    passed = auditor.audit_all_files()
    
    sys.exit(0 if passed else 1)

if __name__ == '__main__':
    main()
