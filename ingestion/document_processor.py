"""
Document Processor for User Uploads
Handles PDF, Excel, PowerPoint, Word, and other document formats
"""

from pathlib import Path
from typing import Dict, Any, Optional, List, BinaryIO
from datetime import datetime
import io
from loguru import logger
from pydantic import BaseModel, Field

# PDF processing
try:
    import PyPDF2
    from pdfminer.high_level import extract_text as pdf_extract_text
except ImportError:
    logger.warning("PDF libraries not fully available")

# Excel processing
import pandas as pd
import openpyxl

# Word processing
from docx import Document as DocxDocument

# PowerPoint processing
try:
    from pptx import Presentation
except ImportError:
    logger.warning("python-pptx not available")


class UploadedDocument(BaseModel):
    """Schema for uploaded documents"""
    document_id: str = Field(..., description="Unique document identifier")
    filename: str = Field(..., description="Original filename")
    file_type: str = Field(..., description="File type/extension")
    upload_date: datetime = Field(default_factory=datetime.utcnow)
    file_size: int = Field(..., description="File size in bytes")
    content: Optional[str] = Field(None, description="Extracted text content")
    metadata: Dict[str, Any] = Field(default_factory=dict, description="Document metadata")
    tables: List[Dict[str, Any]] = Field(default_factory=list, description="Extracted tables")
    processed: bool = Field(False, description="Whether document has been processed")


class DocumentProcessor:
    """Process various document types for data extraction"""
    
    def __init__(self, upload_dir: Optional[Path] = None):
        """
        Initialize document processor
        
        Args:
            upload_dir: Directory for storing uploaded files
        """
        from config.settings import get_settings
        settings = get_settings()
        
        self.upload_dir = upload_dir or (settings.data_dir / "uploads")
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"Document Processor initialized with upload dir: {self.upload_dir}")
    
    def process_upload(
        self,
        file_content: bytes,
        filename: str,
        file_type: Optional[str] = None
    ) -> UploadedDocument:
        """
        Process an uploaded file
        
        Args:
            file_content: File content as bytes
            filename: Original filename
            file_type: File type (auto-detected if not provided)
            
        Returns:
            UploadedDocument object
        """
        if file_type is None:
            file_type = Path(filename).suffix.lower()
        
        logger.info(f"Processing upload: {filename} (type: {file_type})")
        
        # Generate document ID
        doc_id = f"{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}_{filename}"
        
        # Save file
        file_path = self.upload_dir / doc_id
        with open(file_path, 'wb') as f:
            f.write(file_content)
        
        # Process based on file type
        content = None
        tables = []
        metadata = {}
        
        if file_type in ['.pdf']:
            content, tables, metadata = self.process_pdf(file_content)
        elif file_type in ['.xlsx', '.xls', '.xlsm']:
            content, tables, metadata = self.process_excel(file_content, filename)
        elif file_type in ['.docx', '.doc']:
            content, tables, metadata = self.process_word(file_content)
        elif file_type in ['.pptx', '.ppt']:
            content, tables, metadata = self.process_powerpoint(file_content)
        elif file_type in ['.csv']:
            content, tables, metadata = self.process_csv(file_content)
        elif file_type in ['.txt', '.md']:
            content = file_content.decode('utf-8', errors='ignore')
        else:
            logger.warning(f"Unsupported file type: {file_type}")
            content = "Binary file - content extraction not supported"
        
        document = UploadedDocument(
            document_id=doc_id,
            filename=filename,
            file_type=file_type,
            file_size=len(file_content),
            content=content,
            metadata=metadata,
            tables=tables,
            processed=True
        )
        
        logger.info(f"Processed document: {doc_id}")
        return document
    
    def process_pdf(self, file_content: bytes) -> tuple[str, List[Dict], Dict]:
        """Extract text and tables from PDF"""
        try:
            # Try pdfminer for better text extraction
            text = pdf_extract_text(io.BytesIO(file_content))
            
            # Extract metadata using PyPDF2
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            metadata = {
                'num_pages': len(pdf_reader.pages),
                'pdf_info': dict(pdf_reader.metadata) if pdf_reader.metadata else {}
            }
            
            # TODO: Implement table extraction (could use camelot or tabula)
            tables = []
            
            return text, tables, metadata
            
        except Exception as e:
            logger.error(f"PDF processing error: {str(e)}")
            return f"Error processing PDF: {str(e)}", [], {}
    
    def process_excel(self, file_content: bytes, filename: str) -> tuple[str, List[Dict], Dict]:
        """Extract data from Excel files"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(io.BytesIO(file_content))
            
            text_parts = []
            tables = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert to text
                text_parts.append(f"Sheet: {sheet_name}\n")
                text_parts.append(df.to_string())
                text_parts.append("\n\n")
                
                # Store as structured table
                tables.append({
                    'sheet_name': sheet_name,
                    'data': df.to_dict('records'),
                    'columns': df.columns.tolist(),
                    'shape': df.shape
                })
            
            metadata = {
                'num_sheets': len(excel_file.sheet_names),
                'sheet_names': excel_file.sheet_names
            }
            
            text = ''.join(text_parts)
            
            return text, tables, metadata
            
        except Exception as e:
            logger.error(f"Excel processing error: {str(e)}")
            return f"Error processing Excel: {str(e)}", [], {}
    
    def process_word(self, file_content: bytes) -> tuple[str, List[Dict], Dict]:
        """Extract text and tables from Word documents"""
        try:
            doc = DocxDocument(io.BytesIO(file_content))
            
            # Extract paragraphs
            text_parts = []
            for para in doc.paragraphs:
                text_parts.append(para.text)
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    'data': table_data,
                    'num_rows': len(table_data),
                    'num_cols': len(table_data[0]) if table_data else 0
                })
            
            metadata = {
                'num_paragraphs': len(doc.paragraphs),
                'num_tables': len(doc.tables)
            }
            
            text = '\n'.join(text_parts)
            
            return text, tables, metadata
            
        except Exception as e:
            logger.error(f"Word processing error: {str(e)}")
            return f"Error processing Word document: {str(e)}", [], {}
    
    def process_powerpoint(self, file_content: bytes) -> tuple[str, List[Dict], Dict]:
        """Extract text from PowerPoint presentations"""
        try:
            prs = Presentation(io.BytesIO(file_content))
            
            text_parts = []
            tables = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n--- Slide {slide_num} ---\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
                    
                    # Extract tables if present
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        
                        tables.append({
                            'slide_num': slide_num,
                            'data': table_data
                        })
            
            metadata = {
                'num_slides': len(prs.slides)
            }
            
            text = '\n'.join(text_parts)
            
            return text, tables, metadata
            
        except Exception as e:
            logger.error(f"PowerPoint processing error: {str(e)}")
            return f"Error processing PowerPoint: {str(e)}", [], {}
    
    def process_csv(self, file_content: bytes) -> tuple[str, List[Dict], Dict]:
        """Process CSV files"""
        try:
            df = pd.read_csv(io.BytesIO(file_content))
            
            text = df.to_string()
            
            tables = [{
                'data': df.to_dict('records'),
                'columns': df.columns.tolist(),
                'shape': df.shape
            }]
            
            metadata = {
                'num_rows': len(df),
                'num_columns': len(df.columns),
                'columns': df.columns.tolist()
            }
            
            return text, tables, metadata
            
        except Exception as e:
            logger.error(f"CSV processing error: {str(e)}")
            return f"Error processing CSV: {str(e)}", [], {}
    
    def extract_financial_data(self, document: UploadedDocument) -> Dict[str, Any]:
        """
        Extract financial data from uploaded documents using LLM
        
        Args:
            document: Uploaded document
            
        Returns:
            Extracted financial data
        """
        from utils import LLMClient
        
        llm = LLMClient()
        
        prompt = f"""Extract financial data from the following document content.
Look for:
- Revenue figures
- EBITDA, EBIT, Net Income
- Balance sheet items (Assets, Liabilities, Equity)
- Cash flows
- Key metrics and ratios
- Dates and periods

Document: {document.filename}
Content:
{document.content[:8000]}

Provide response as structured JSON."""

        messages = [
            {"role": "system", "content": "You are a financial analyst extracting structured data from documents."},
            {"role": "user", "content": prompt}
        ]
        
        try:
            response = llm.chat(messages, temperature=0.1)
            import json
            return json.loads(response)
        except Exception as e:
            logger.error(f"Error extracting financial data: {str(e)}")
            return {"error": str(e), "raw_response": response if 'response' in locals() else None}
